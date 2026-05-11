"""
Build the final slide deck for the Diabetes Risk Prediction project.
Source of truth: outputs/final_report.md and outputs/results_packet.md.
All numbers below are taken from the validated R run (runtime_output_full.md).
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

OUT_DIR = Path(r"C:\Users\JSEer\AEI-POC\outputs")
TREE_PNG = OUT_DIR / "decision_tree_plot.png"
DECK_PATH = OUT_DIR / "final_slide_deck.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W, SLIDE_H = prs.slide_width, prs.slide_height

NAVY = RGBColor(0x0B, 0x2A, 0x4A)
ACCENT = RGBColor(0xC1, 0x44, 0x1E)
GREY = RGBColor(0x55, 0x55, 0x55)
BLACK = RGBColor(0x10, 0x10, 0x10)


def add_blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_title_bar(slide, text):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), SLIDE_W - Inches(1.0), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = NAVY


def add_textbox(slide, left, top, width, height, lines, size=18, bold=False, color=BLACK):
    """lines may be a list of (text, indent, bold, size, color) tuples or plain strings."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, indent, b, sz, col = item, 0, bold, size, color
        else:
            text, indent, b, sz, col = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = indent
        for run in p.runs:
            run.font.size = Pt(sz)
            run.font.bold = b
            run.font.color.rgb = col


def add_footer(slide, n_total, n_current):
    box = slide.shapes.add_textbox(Inches(0.4), SLIDE_H - Inches(0.45), SLIDE_W - Inches(0.8), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Diabetes Risk Prediction by Patient Health Indicators  |  J. Ermert  |  MIS 401  |  Slide {n_current} of {n_total}"
    p.alignment = PP_ALIGN.RIGHT
    for r in p.runs:
        r.font.size = Pt(10)
        r.font.color.rgb = GREY


# --------- Slide content ---------

slides_spec = []

# Slide 1 - Title
slides_spec.append(("title", {}))

# Slide 2 - Why this matters
slides_spec.append(("bullets", {
    "title": "Why This Matters",
    "bullets": [
        ("Diabetes is a leading U.S. public-health condition with high cost and uneven distribution.", 0),
        ("Identifying which patient indicators predict diabetes supports earlier screening and targeted prevention.", 0),
        ("Question: how well do classification methods separate Diabetes / Prediabetes / NoDiabetes — and which indicators carry the signal?", 0),
        ("Personal motivation: Type II diabetes in family + prior biopharma BI/AI analytics work.", 0),
    ],
}))

# Slide 3 - Dataset & response variable
slides_spec.append(("bullets", {
    "title": "Dataset and Response Variable",
    "bullets": [
        ("Source: CDC Diabetes Health Indicators (UCI ML Repository, BRFSS-derived).", 0),
        ("253,680 observations  |  21 predictors + Diabetes_012  |  No missing values.", 0),
        ("Response Diabetes_012 has three classes:", 0),
        ("0 = NoDiabetes — 213,703 (84.24 %)", 1),
        ("1 = Prediabetes — 4,631 (1.83 %)", 1),
        ("2 = Diabetes — 35,346 (13.93 %)", 1),
        ("Severe class imbalance is the single most important property of the dataset.", 0),
    ],
}))

# Slide 4 - Variables
slides_spec.append(("two_col", {
    "title": "Variables Used",
    "left_header": "Categorical / binary",
    "left_bullets": [
        "HighBP", "HighChol", "CholCheck", "Smoker", "Stroke",
        "HeartDiseaseorAttack", "PhysActivity", "Fruits", "Veggies",
        "HvyAlcoholConsump", "AnyHealthcare", "NoDocbcCost",
        "DiffWalk", "Sex",
    ],
    "right_header": "Continuous / ordinal",
    "right_bullets": [
        "BMI",
        "GenHlth (1 excellent – 5 poor)",
        "MentHlth (days in past 30)",
        "PhysHlth (days in past 30)",
        "Age (1 youngest – 13 oldest)",
        "Education",
        "Income",
    ],
    "footnote": "Categorical predictors converted to factors for logit + tree; KNN used numeric encoding with z-score scaling.",
}))

# Slide 5 - Preliminary analysis
slides_spec.append(("bullets", {
    "title": "Preliminary Analysis (EDA Highlights)",
    "bullets": [
        ("BMI shifts upward with diabetes status (median ~27 → 30 → 31).", 0),
        ("Age category increases with diabetes status (median 8 → 9 → 10).", 0),
        ("Self-rated General Health is markedly worse in the diabetes group.", 0),
        ("HighBP is roughly 3× more common in the diabetes group than in NoDiabetes.", 0),
        ("HighChol shows the same direction at slightly smaller magnitude.", 0),
        ("PhysActivity is somewhat lower; Sex effect present but weaker than the four above.", 0),
        ("Strongest early signals: HighBP, GenHlth, HighChol, BMI, Age.", 0),
    ],
}))

# Slide 6 - Methods
slides_spec.append(("bullets", {
    "title": "Methods",
    "bullets": [
        ("Same 70/30 train/test split with set.seed(401) for all three models.", 0),
        ("Train n = 177,576  |  Test n = 76,104.", 0),
        ("Multinomial Logistic Regression — nnet::multinom on factor-coded predictors.", 0),
        ("KNN — class::knn on z-scored numeric predictors; tuned over k ∈ {3, 5, 7, 9}.", 0),
        ("Decision Tree — rpart with method = 'class'.", 0),
        ("Default rpart (cp = 0.01) refused to split — tuned to cp = 0.001, minsplit = 2000, minbucket = 1000, maxdepth = 5 to grow a usable shallow tree.", 1),
        ("Evaluation: confusion matrix + overall test accuracy + tree variable importance.", 0),
    ],
}))

# Slide 7 - Model comparison
slides_spec.append(("table", {
    "title": "Model Comparison (Test Accuracy)",
    "headers": ["Rank", "Model", "Test Accuracy"],
    "rows": [
        ["1", "Decision Tree (cp = 0.001, depth = 5)", "0.8475"],
        ["2", "Multinomial Logistic Regression", "0.8467"],
        ["3", "KNN (best K = 9)", "0.8382"],
        ["4", "KNN (k = 5)", "0.8301"],
    ],
    "callouts": [
        "All four variants land within ~2 points of each other — and within ~0.5 points of the 84.2 % majority-class baseline.",
        "Accuracy alone is misleading on this dataset; per-class evidence follows.",
    ],
}))

# Slide 8 - Decision tree image
slides_spec.append(("image", {
    "title": "Decision Tree — Structure",
    "image": str(TREE_PNG),
    "caption": "One Diabetes-predicting leaf: HighBP=1 ∧ GenHlth ≥ 4 ∧ BMI ≥ 28 ∧ HighChol=1 ∧ BMI ≥ 35 (n = 3,961, ~60 % Diabetes).",
}))

# Slide 9 - Variable importance & key insights
slides_spec.append(("bullets", {
    "title": "Variable Importance and Key Insights",
    "bullets": [
        ("Decision-tree variable importance (top): HighBP > GenHlth > HighChol > DiffWalk > Age > PhysHlth > BMI > HeartDiseaseorAttack.", 0),
        ("Lines up directly with EDA: HighBP, GenHlth, HighChol, BMI, Age.", 1),
        ("Aligns with established Type II diabetes risk-factor literature.", 1),
        ("All models predict NoDiabetes for ~98 % of test rows; only KNN k = 5 ever predicts Prediabetes (88 times, 2 correct).", 0),
        ("Diabetes-class sensitivity: Tree 9 %  |  Logit 17 %  |  KNN k=5 22 %.", 0),
        ("Decision Tree wins on accuracy by trading sensitivity for precision; logit is the most balanced single deployable model.", 0),
    ],
}))

# Slide 10 - Practical value & conclusion
slides_spec.append(("bullets", {
    "title": "Practical Value and Conclusion",
    "bullets": [
        ("Even at modest sensitivity, the tree is operationally useful in two ways:", 0),
        ("\"Rule-out\" branch: 101,437 patients with HighBP = 0 → only 6 % Diabetes rate. Safe to deprioritize from intensive screening.", 1),
        ("\"High-risk\" leaf: 3,961 patients (BMI ≥ 35 + HighChol + GenHlth poor + HighBP = 1) → ~60 % Diabetes rate, ~5× the population base rate.", 1),
        ("Honest framing: accuracy ≈ baseline; class imbalance dominates results.", 0),
        ("Clearest improvement path is methodological — class weights, oversampling, or a cost-sensitive threshold — not a different algorithm.", 0),
        ("AI disclosure: generative AI organized + polished structure; data, methods, and analytical decisions are the author's, supported by validated R outputs.", 0),
    ],
}))

# --------- Render ---------

n_total = len(slides_spec)

for idx, (kind, spec) in enumerate(slides_spec, start=1):
    slide = add_blank_slide()

    if kind == "title":
        # Title slide
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), SLIDE_W - Inches(1.6), Inches(2.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Diabetes Risk Prediction by Patient Health Indicators"
        p.alignment = PP_ALIGN.LEFT
        for r in p.runs:
            r.font.size = Pt(40)
            r.font.bold = True
            r.font.color.rgb = NAVY
        p2 = tf.add_paragraph()
        p2.text = "UCI x CDC  |  MIS 401: Business Intelligence & Analytics"
        for r in p2.runs:
            r.font.size = Pt(22)
            r.font.color.rgb = ACCENT

        sub = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), SLIDE_W - Inches(1.6), Inches(1.5))
        sf = sub.text_frame
        sf.word_wrap = True
        for i, line in enumerate([
            "Joshua Ermert",
            "Dr. Xialu Liu",
            "7 May 2026",
        ]):
            p = sf.paragraphs[0] if i == 0 else sf.add_paragraph()
            p.text = line
            for r in p.runs:
                r.font.size = Pt(18)
                r.font.color.rgb = BLACK

    elif kind == "bullets":
        add_title_bar(slide, spec["title"])
        body = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.4), SLIDE_W - Inches(1.2), SLIDE_H - Inches(2.0)
        )
        tf = body.text_frame
        tf.word_wrap = True
        for i, item in enumerate(spec["bullets"]):
            text, indent = item if isinstance(item, tuple) else (item, 0)
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = ("• " if indent == 0 else "– ") + text
            p.level = indent
            for r in p.runs:
                r.font.size = Pt(20 if indent == 0 else 17)
                r.font.color.rgb = BLACK

    elif kind == "two_col":
        add_title_bar(slide, spec["title"])
        col_w = (SLIDE_W - Inches(1.4)) / 2

        for col_idx, (header, items) in enumerate([
            (spec["left_header"], spec["left_bullets"]),
            (spec["right_header"], spec["right_bullets"]),
        ]):
            left = Inches(0.6) + col_idx * (col_w + Inches(0.2))
            box = slide.shapes.add_textbox(left, Inches(1.4), col_w, Inches(5.2))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = header
            for r in p.runs:
                r.font.size = Pt(22)
                r.font.bold = True
                r.font.color.rgb = ACCENT
            for it in items:
                p = tf.add_paragraph()
                p.text = "• " + it
                for r in p.runs:
                    r.font.size = Pt(16)
                    r.font.color.rgb = BLACK

        if spec.get("footnote"):
            fn = slide.shapes.add_textbox(Inches(0.6), Inches(6.7), SLIDE_W - Inches(1.2), Inches(0.5))
            ftf = fn.text_frame
            ftf.word_wrap = True
            p = ftf.paragraphs[0]
            p.text = spec["footnote"]
            for r in p.runs:
                r.font.size = Pt(13)
                r.font.italic = True
                r.font.color.rgb = GREY

    elif kind == "table":
        add_title_bar(slide, spec["title"])
        rows_n = len(spec["rows"]) + 1
        cols_n = len(spec["headers"])
        tbl_left = Inches(0.8)
        tbl_top = Inches(1.5)
        tbl_w = SLIDE_W - Inches(1.6)
        tbl_h = Inches(0.5 * rows_n + 0.4)
        tbl_shape = slide.shapes.add_table(rows_n, cols_n, tbl_left, tbl_top, tbl_w, tbl_h)
        tbl = tbl_shape.table

        for c, h in enumerate(spec["headers"]):
            cell = tbl.cell(0, c)
            cell.text = h
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(18)
                    r.font.bold = True
                    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY

        for r_idx, row in enumerate(spec["rows"], start=1):
            for c_idx, val in enumerate(row):
                cell = tbl.cell(r_idx, c_idx)
                cell.text = val
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(16)
                        run.font.color.rgb = BLACK

        if spec.get("callouts"):
            cb = slide.shapes.add_textbox(Inches(0.8), tbl_top + tbl_h + Inches(0.3),
                                          SLIDE_W - Inches(1.6), Inches(2.0))
            ctf = cb.text_frame
            ctf.word_wrap = True
            for i, c in enumerate(spec["callouts"]):
                p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
                p.text = "• " + c
                for run in p.runs:
                    run.font.size = Pt(16)
                    run.font.color.rgb = ACCENT if i == 0 else BLACK

    elif kind == "image":
        add_title_bar(slide, spec["title"])
        img_w = SLIDE_W - Inches(2.0)
        slide.shapes.add_picture(spec["image"], Inches(1.0), Inches(1.4), width=img_w)
        cap = slide.shapes.add_textbox(Inches(0.6), SLIDE_H - Inches(1.0),
                                       SLIDE_W - Inches(1.2), Inches(0.6))
        ctf = cap.text_frame
        ctf.word_wrap = True
        p = ctf.paragraphs[0]
        p.text = spec["caption"]
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(14)
            r.font.italic = True
            r.font.color.rgb = GREY

    add_footer(slide, n_total, idx)

prs.save(DECK_PATH)
print(f"Saved {DECK_PATH}")
