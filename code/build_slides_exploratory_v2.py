"""
Exploratory v2 deck — same 10-slide structure, frozen numbers, frozen ranking.
Elevations over final_slide_deck_polished.pptx:
  - section tags top-right ("02 — WHY THIS MATTERS" style)
  - signature dot-grid class imbalance chart on slide 3
  - magnitude / direction indicators on slide 5 EDA cards
  - refined model-comparison chart with gap-from-baseline annotations (slide 7)
  - EDA-callback variable-importance chart + refined sensitivity chart (slide 9)
  - proportion badges on slide 10 ("~40% of patients", "~2% of patients")
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

OUT = Path(r"C:\Users\JSEer\AEI-POC\outputs")
DECK = OUT / "final_slide_deck_exploratory_v2.pptx"

NAVY = RGBColor(0x0B, 0x2A, 0x4A)
RUST = RGBColor(0xC1, 0x44, 0x1E)
TEAL = RGBColor(0x1F, 0x7A, 0x8C)
AMBER = RGBColor(0xE0, 0xA4, 0x58)
GREY_DARK = RGBColor(0x3A, 0x3A, 0x3A)
GREY_MID = RGBColor(0x77, 0x77, 0x77)
GREY_LINE = RGBColor(0xC9, 0xCD, 0xD2)
GREY_BG = RGBColor(0xF4, 0xF6, 0xF8)
RUST_TINT = RGBColor(0xF7, 0xE7, 0xDF)
TEAL_TINT = RGBColor(0xE0, 0xEE, 0xF1)
NAVY_TINT = RGBColor(0xE2, 0xE7, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

FONT = "Calibri"
TOTAL = 10


# ----------------------------------------------------------------------- helpers

def add_blank():
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def set_line(shape, color, width_pt=1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def set_text(tf, text, *, size=18, bold=False, color=GREY_DARK,
             align=PP_ALIGN.LEFT, italic=False, font=FONT):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def add_paragraph(tf, text, *, size=18, bold=False, color=GREY_DARK,
                  align=PP_ALIGN.LEFT, italic=False, level=0, font=FONT,
                  space_before=2):
    p = tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def add_textbox(slide, left, top, width, height, text="",
                size=18, bold=False, color=GREY_DARK, align=PP_ALIGN.LEFT,
                italic=False, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    if text:
        set_text(tf, text, size=size, bold=bold, color=color,
                 align=align, italic=italic)
    return tb, tf


def add_rounded_card(slide, left, top, width, height, *,
                     fill=GREY_BG, border=None, border_w=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, top, width, height)
    shp.adjustments[0] = 0.05
    fill_solid(shp, fill)
    if border is None:
        no_line(shp)
    else:
        set_line(shp, border, border_w)
    shp.shadow.inherit = False
    return shp


def add_accent_line(slide, left, top, width, color=RUST, height=Pt(3.5)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill_solid(bar, color)
    no_line(bar)
    return bar


def add_header_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 0, 0, SW, Inches(0.18))
    fill_solid(bar, NAVY)
    no_line(bar)


def add_footer(slide, page_n, page_total):
    add_textbox(
        slide, Inches(0.45), SH - Inches(0.45),
        Inches(8.5), Inches(0.3),
        "Diabetes Risk Prediction by Patient Health Indicators  ·  J. Ermert  ·  MIS 401",
        size=10, color=GREY_MID, align=PP_ALIGN.LEFT,
    )
    add_textbox(
        slide, SW - Inches(2.5), SH - Inches(0.45),
        Inches(2.0), Inches(0.3),
        f"{page_n} / {page_total}",
        size=10, color=GREY_MID, align=PP_ALIGN.RIGHT, bold=True,
    )


def add_section_tag(slide, n, name):
    """Section indicator under header bar, top-right corner."""
    add_textbox(
        slide, SW - Inches(6.5), Inches(0.30),
        Inches(6.0), Inches(0.30),
        f"{n:02d}  —  {name}",
        size=10, bold=True, color=GREY_MID, align=PP_ALIGN.RIGHT,
    )


def add_slide_title(slide, title_text, section_n=None, section_name=None):
    add_header_bar(slide)
    if section_n is not None and section_name:
        add_section_tag(slide, section_n, section_name)
    add_textbox(
        slide, Inches(0.5), Inches(0.45), SW - Inches(1.0), Inches(0.85),
        title_text, size=30, bold=True, color=NAVY,
    )
    add_accent_line(slide, Inches(0.5), Inches(1.30), Inches(0.7), color=RUST)


# ----------------------------------------------------------------------- slides


def slide_1_title():
    s = add_blank()
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(4.4), SH)
    fill_solid(band, NAVY)
    no_line(band)

    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0.6), Inches(3.6),
                                Inches(0.85), Pt(4))
    fill_solid(accent, RUST)
    no_line(accent)

    add_textbox(s, Inches(0.6), Inches(2.9), Inches(3.6), Inches(0.5),
                "MIS 401  ·  BUSINESS INTELLIGENCE & ANALYTICS",
                size=12, bold=True, color=WHITE)

    add_textbox(s, Inches(0.6), Inches(3.9), Inches(3.6), Inches(2.0),
                "FINAL\nPROJECT",
                size=44, bold=True, color=WHITE)

    add_textbox(s, Inches(0.6), Inches(6.4), Inches(3.6), Inches(0.5),
                "UCI x CDC", size=14, bold=True, color=AMBER)

    add_textbox(s, Inches(5.0), Inches(2.4), Inches(7.8), Inches(2.4),
                "Diabetes Risk Prediction\nby Patient Health Indicators",
                size=42, bold=True, color=NAVY)

    add_accent_line(s, Inches(5.0), Inches(4.55), Inches(1.2), color=RUST)

    add_textbox(s, Inches(5.0), Inches(4.85), Inches(7.8), Inches(0.5),
                "Comparing logistic regression, KNN, and a tuned decision tree on 253,680 patients.",
                size=16, italic=True, color=GREY_DARK)

    add_textbox(s, Inches(5.0), Inches(6.0), Inches(7.8), Inches(0.4),
                "Joshua Ermert", size=18, bold=True, color=NAVY)
    add_textbox(s, Inches(5.0), Inches(6.4), Inches(7.8), Inches(0.4),
                "Dr. Xialu Liu  ·  7 May 2026",
                size=13, color=GREY_MID)


def slide_2_why():
    s = add_blank()
    add_slide_title(s, "Why This Matters", 2, "CONTEXT")

    quote = add_rounded_card(s, Inches(0.5), Inches(1.7), Inches(5.6), Inches(4.7),
                             fill=NAVY)
    no_line(quote)
    add_textbox(s, Inches(0.85), Inches(1.95), Inches(5.0), Inches(1.0),
                "“", size=72, bold=True, color=AMBER)
    add_textbox(s, Inches(0.85), Inches(2.95), Inches(5.0), Inches(2.6),
                "Identify which patient indicators predict diabetes — and turn that signal into earlier screening and targeted prevention.",
                size=20, color=WHITE)
    add_accent_line(s, Inches(0.85), Inches(5.55), Inches(0.8), color=RUST)
    add_textbox(s, Inches(0.85), Inches(5.7), Inches(5.0), Inches(0.6),
                "The driving question for this project.",
                size=13, italic=True, color=NAVY_TINT)

    headers = [
        ("Public-health weight",
         "Diabetes is a leading U.S. public-health condition with high cost and uneven distribution."),
        ("Three classifiers, one comparison",
         "Multinomial logistic regression, KNN, and a tuned decision tree applied to the same 70/30 split."),
        ("Personal stake",
         "Type II diabetes in the family, plus prior biopharma BI/AI analytics work — applied analytics in a real healthcare context."),
    ]
    top = Inches(1.7)
    for title, body in headers:
        card = add_rounded_card(s, Inches(6.5), top, Inches(6.4), Inches(1.45),
                                fill=GREY_BG)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(6.5), top + Inches(0.15),
                                 Pt(4), Inches(1.15))
        fill_solid(bar, RUST)
        no_line(bar)
        add_textbox(s, Inches(6.7), top + Inches(0.15), Inches(6.0), Inches(0.4),
                    title, size=15, bold=True, color=NAVY)
        add_textbox(s, Inches(6.7), top + Inches(0.55), Inches(6.0), Inches(0.85),
                    body, size=13, color=GREY_DARK)
        top += Inches(1.6)

    add_footer(s, 2, TOTAL)


def slide_3_dataset():
    s = add_blank()
    add_slide_title(s, "Dataset and Response Variable", 3, "DATA")

    facts = [
        ("Source", "CDC Diabetes Health Indicators (UCI ML Repository, BRFSS-derived)"),
        ("Observations", "253,680"),
        ("Predictors", "21  +  1 response (Diabetes_012)"),
        ("Missing values", "None"),
    ]
    top = Inches(1.7)
    for label, val in facts:
        add_textbox(s, Inches(0.5), top, Inches(2.6), Inches(0.35),
                    label.upper(), size=11, bold=True, color=GREY_MID)
        add_textbox(s, Inches(0.5), top + Inches(0.32), Inches(5.5), Inches(0.55),
                    val, size=15, bold=True, color=NAVY)
        top += Inches(0.95)

    # Right: signature dot-grid v2 chart
    s.shapes.add_picture(str(OUT / "chart_class_imbalance_v2.png"),
                         Inches(6.2), Inches(1.6),
                         width=Inches(6.9))

    callout = add_rounded_card(s, Inches(0.5), Inches(5.6),
                               Inches(5.5), Inches(1.3),
                               fill=RUST_TINT, border=RUST, border_w=1.25)
    add_textbox(s, Inches(0.75), Inches(5.7), Inches(5.0), Inches(0.4),
                "KEY TAKEAWAY", size=11, bold=True, color=RUST)
    add_textbox(s, Inches(0.75), Inches(6.0), Inches(5.0), Inches(0.85),
                "84.24 / 1.83 / 13.93 — class imbalance is the single most important property of this dataset.",
                size=12, color=GREY_DARK)

    add_footer(s, 3, TOTAL)


def slide_4_variables():
    s = add_blank()
    add_slide_title(s, "Variables Used", 4, "PREDICTORS")

    cat = ["HighBP", "HighChol", "CholCheck", "Smoker", "Stroke",
           "HeartDiseaseorAttack", "PhysActivity", "Fruits", "Veggies",
           "HvyAlcoholConsump", "AnyHealthcare", "NoDocbcCost",
           "DiffWalk", "Sex"]
    cont = ["BMI", "GenHlth (1 excellent – 5 poor)",
            "MentHlth (days in past 30)", "PhysHlth (days in past 30)",
            "Age (1 youngest – 13 oldest)", "Education", "Income"]

    def col(left, header, items, header_color):
        add_rounded_card(s, left, Inches(1.7),
                         Inches(6.0), Inches(4.8),
                         fill=GREY_BG)
        add_textbox(s, left + Inches(0.35), Inches(1.95),
                    Inches(5.4), Inches(0.4),
                    header.upper(), size=11, bold=True, color=GREY_MID)
        add_textbox(s, left + Inches(0.35), Inches(2.25),
                    Inches(5.4), Inches(0.55),
                    f"{len(items)} predictors", size=22, bold=True,
                    color=header_color)

        col_w = Inches(2.6)
        gap = Inches(0.2)
        per_col = (len(items) + 1) // 2
        for c in range(2):
            tb = s.shapes.add_textbox(
                left + Inches(0.35) + c * (col_w + gap),
                Inches(3.0), col_w, Inches(3.5))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = Pt(0)
            tf.margin_top = Pt(0)
            chunk = items[c * per_col:(c + 1) * per_col]
            for i, item in enumerate(chunk):
                if i == 0:
                    set_text(tf, "•  " + item,
                             size=13, color=GREY_DARK)
                else:
                    add_paragraph(tf, "•  " + item,
                                  size=13, color=GREY_DARK,
                                  space_before=4)

    col(Inches(0.5), "Categorical / binary", cat, NAVY)
    col(Inches(6.85), "Continuous / ordinal", cont, TEAL)

    add_textbox(s, Inches(0.5), Inches(6.65), Inches(12.4), Inches(0.4),
                "Categorical predictors → factors for logit + tree.   KNN → numeric encoding with z-score scaling fit on the training set.",
                size=12, italic=True, color=GREY_MID)

    add_footer(s, 4, TOTAL)


def slide_5_eda():
    s = add_blank()
    add_slide_title(s, "Preliminary Analysis  —  EDA highlights", 5, "EDA")

    # Each finding now carries: variable name, body, direction arrow, magnitude (1-3 dots)
    findings = [
        ("BMI",          "Median rises 27  →  30  →  31",                   "↑", 2),
        ("Age",          "Median age category 8  →  9  →  10",               "↑", 2),
        ("GenHlth",      "Markedly worse self-rated health in Diabetes group","↑", 3),
        ("HighBP",       "~3× more common in Diabetes group than NoDiabetes","↑", 3),
        ("HighChol",     "Same direction at slightly smaller magnitude",     "↑", 2),
        ("PhysActivity", "Lower in Diabetes group; Sex effect weaker",       "↓", 1),
    ]
    grid_left = Inches(0.5)
    grid_top = Inches(1.7)
    cell_w = Inches(4.05)
    cell_h = Inches(1.55)
    gap_x = Inches(0.2)
    gap_y = Inches(0.2)

    for idx, (label, body, arrow, mag) in enumerate(findings):
        r, c = divmod(idx, 3)
        x = grid_left + c * (cell_w + gap_x)
        y = grid_top + r * (cell_h + gap_y)
        add_rounded_card(s, x, y, cell_w, cell_h, fill=GREY_BG)
        add_accent_line(s, x + Inches(0.3), y + Inches(0.25),
                        Inches(0.5), color=RUST, height=Pt(3))

        # Variable name (left) + direction arrow (right)
        add_textbox(s, x + Inches(0.3), y + Inches(0.35),
                    Inches(2.4), Inches(0.5),
                    label, size=18, bold=True, color=NAVY)
        arrow_color = RUST if arrow == "↑" else TEAL
        add_textbox(s, x + cell_w - Inches(1.2), y + Inches(0.30),
                    Inches(0.9), Inches(0.55),
                    arrow, size=24, bold=True, color=arrow_color,
                    align=PP_ALIGN.RIGHT)

        # Body
        add_textbox(s, x + Inches(0.3), y + Inches(0.85),
                    cell_w - Inches(0.6), Inches(0.55),
                    body, size=12, color=GREY_DARK)

        # Magnitude indicator: three dots, filled = mag
        dot_y = y + Inches(1.30)
        for d in range(3):
            dot = s.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x + Inches(0.3) + Inches(0.18) * d,
                dot_y, Inches(0.13), Inches(0.13),
            )
            fill_solid(dot, RUST if d < mag else GREY_LINE)
            no_line(dot)
        add_textbox(s, x + Inches(1.0), dot_y - Inches(0.02),
                    cell_w - Inches(1.2), Inches(0.18),
                    "magnitude", size=8, bold=True, color=GREY_MID)

    callout = add_rounded_card(s, Inches(0.5), Inches(5.3),
                               Inches(12.4), Inches(0.95),
                               fill=NAVY)
    add_textbox(s, Inches(0.85), Inches(5.4), Inches(2.7), Inches(0.3),
                "STRONGEST EARLY SIGNALS",
                size=11, bold=True, color=AMBER)
    add_textbox(s, Inches(0.85), Inches(5.65), Inches(11.5), Inches(0.6),
                "HighBP   ·   GenHlth   ·   HighChol   ·   BMI   ·   Age",
                size=20, bold=True, color=WHITE)

    add_textbox(s, Inches(0.5), Inches(6.45), Inches(12.4), Inches(0.4),
                "These five predictors should reappear in the modeling output. Hold onto them — they show up again on slide 9.",
                size=12, italic=True, color=GREY_MID)

    add_footer(s, 5, TOTAL)


def slide_6_methods():
    s = add_blank()
    add_slide_title(s, "Methods", 6, "MODELS")

    add_textbox(s, Inches(0.5), Inches(1.6), Inches(12.4), Inches(0.4),
                "Same split for all three models  ·  set.seed(401)  ·  70/30  ·  Train n = 177,576   Test n = 76,104",
                size=14, italic=True, color=GREY_MID)

    methods = [
        {
            "color": NAVY,
            "tag": "STATISTICAL",
            "name": "Multinomial\nLogistic Regression",
            "pkg": "nnet::multinom",
            "note": "Factor-coded predictors. Interpretable coefficients per non-reference class (Prediabetes, Diabetes vs. NoDiabetes).",
        },
        {
            "color": TEAL,
            "tag": "INSTANCE-BASED",
            "name": "K-Nearest\nNeighbors",
            "pkg": "class::knn",
            "note": "z-score scaling fit on training set. Tuned k ∈ {3, 5, 7, 9}; reports best K.",
        },
        {
            "color": RUST,
            "tag": "RULE-BASED",
            "name": "Decision Tree",
            "pkg": "rpart  (method = \"class\")",
            "note": "Default cp = 0.01 refused to split. Tuned to cp = 0.001, minsplit = 2000, minbucket = 1000, maxdepth = 5 — honest tuning, not deep optimization.",
        },
    ]

    card_w = Inches(4.05)
    card_h = Inches(4.4)
    card_top = Inches(2.15)
    gap = Inches(0.2)
    left = Inches(0.5)
    for m in methods:
        add_rounded_card(s, left, card_top, card_w, card_h, fill=WHITE,
                         border=GREY_LINE, border_w=1.0)
        band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  left, card_top, card_w, Inches(0.55))
        fill_solid(band, m["color"])
        no_line(band)
        add_textbox(s, left + Inches(0.3), card_top + Inches(0.12),
                    card_w - Inches(0.6), Inches(0.35),
                    m["tag"], size=11, bold=True, color=WHITE)

        add_textbox(s, left + Inches(0.3), card_top + Inches(0.75),
                    card_w - Inches(0.6), Inches(1.05),
                    m["name"], size=22, bold=True, color=NAVY)

        add_rounded_card(s,
                         left + Inches(0.3),
                         card_top + Inches(2.05),
                         card_w - Inches(0.6), Inches(0.45),
                         fill=GREY_BG)
        add_textbox(s,
                    left + Inches(0.3), card_top + Inches(2.10),
                    card_w - Inches(0.6), Inches(0.35),
                    m["pkg"], size=12, bold=True, color=GREY_DARK,
                    align=PP_ALIGN.CENTER)

        add_textbox(s, left + Inches(0.3), card_top + Inches(2.7),
                    card_w - Inches(0.6), card_h - Inches(2.85),
                    m["note"], size=12, color=GREY_DARK)

        left += card_w + gap

    add_textbox(s, Inches(0.5), Inches(6.7), Inches(12.4), Inches(0.4),
                "Evaluation: confusion matrix + overall test accuracy + decision-tree variable importance.",
                size=12, italic=True, color=GREY_MID)

    add_footer(s, 6, TOTAL)


def slide_7_comparison():
    s = add_blank()
    add_slide_title(s, "Model Comparison  —  Test accuracy", 7, "RESULTS")

    # v2 chart already carries gap-from-baseline pp annotations next to each bar
    s.shapes.add_picture(str(OUT / "chart_model_comparison_v2.png"),
                         Inches(0.5), Inches(1.6), width=Inches(8.6))

    # Compact right column — chart now communicates the gap, so panel is leaner
    add_rounded_card(s, Inches(9.4), Inches(1.7),
                     Inches(3.5), Inches(4.6),
                     fill=NAVY)
    add_textbox(s, Inches(9.65), Inches(1.9), Inches(3.0), Inches(0.4),
                "HEADLINE", size=11, bold=True, color=AMBER)
    add_textbox(s, Inches(9.65), Inches(2.2), Inches(3.0), Inches(1.0),
                "0.8475", size=44, bold=True, color=WHITE)
    add_textbox(s, Inches(9.65), Inches(3.15), Inches(3.0), Inches(0.4),
                "Decision Tree wins —", size=14, bold=True, color=AMBER)
    add_textbox(s, Inches(9.65), Inches(3.4), Inches(3.0), Inches(0.5),
                "by 0.0008 over logit", size=14, color=WHITE)
    add_accent_line(s, Inches(9.65), Inches(4.1), Inches(0.8), color=RUST)

    # Mini gap-from-baseline mini-table inside the panel
    add_textbox(s, Inches(9.65), Inches(4.25), Inches(3.0), Inches(0.35),
                "GAP FROM BASELINE", size=10, bold=True, color=AMBER)
    gap_lines = [
        ("Tree",  "+0.51 pp"),
        ("Logit", "+0.43 pp"),
        ("KNN k=9", "−0.42 pp"),
        ("KNN k=5", "−1.23 pp"),
    ]
    g_top = Inches(4.6)
    for label, val in gap_lines:
        add_textbox(s, Inches(9.65), g_top, Inches(1.7), Inches(0.3),
                    label, size=12, color=WHITE)
        add_textbox(s, Inches(11.3), g_top, Inches(1.4), Inches(0.3),
                    val, size=12, bold=True, color=AMBER,
                    align=PP_ALIGN.RIGHT)
        g_top += Inches(0.36)

    add_textbox(s, Inches(0.5), Inches(6.5), Inches(8.6), Inches(0.4),
                "Accuracy alone is misleading on this dataset — per-class evidence on slide 9.",
                size=13, italic=True, color=GREY_MID)

    add_footer(s, 7, TOTAL)


def slide_8_tree():
    s = add_blank()
    add_slide_title(s, "Decision Tree  —  Structure", 8, "TREE STRUCTURE")

    img_h = Inches(4.7)
    img_w = Inches(4.7 * 1400 / 900)
    img_left = (SW - img_w) / 2
    img_top = Inches(1.55)
    s.shapes.add_picture(str(OUT / "decision_tree_plot.png"),
                         img_left, img_top,
                         width=img_w, height=img_h)

    cap_top = Inches(6.40)
    add_rounded_card(s, Inches(0.5), cap_top,
                     Inches(12.4), Inches(0.6),
                     fill=RUST_TINT)
    add_textbox(s, Inches(0.7), cap_top + Inches(0.05),
                Inches(12.0), Inches(0.5),
                "One Diabetes-predicting leaf:  HighBP=1  ∧  GenHlth ≥ 4  ∧  BMI ≥ 28  ∧  HighChol=1  ∧  BMI ≥ 35   →   n = 3,961, ~60 % Diabetes.",
                size=13, italic=True, color=GREY_DARK, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, 8, TOTAL)


def slide_9_importance():
    s = add_blank()
    add_slide_title(s, "Variable Importance and Per-Class Story", 9, "INTERPRETATION")

    # v2 importance chart (with EDA-callback markers in amber)
    s.shapes.add_picture(str(OUT / "chart_variable_importance_v2.png"),
                         Inches(0.5), Inches(1.65),
                         width=Inches(7.0))

    # v2 sensitivity chart
    s.shapes.add_picture(str(OUT / "chart_diabetes_sensitivity_v2.png"),
                         Inches(7.7), Inches(1.65),
                         width=Inches(5.4))

    callout = add_rounded_card(s, Inches(7.7), Inches(4.55),
                               Inches(5.4), Inches(2.0),
                               fill=GREY_BG, border=GREY_LINE, border_w=1.0)
    add_textbox(s, Inches(7.95), Inches(4.7), Inches(5.0), Inches(0.4),
                "WHAT THIS MEANS", size=11, bold=True, color=RUST)

    bullets_box = s.shapes.add_textbox(Inches(7.95), Inches(5.0),
                                       Inches(5.0), Inches(1.5))
    tf = bullets_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_top = Pt(0)
    set_text(tf, "Tree wins on accuracy by trading sensitivity for precision.",
             size=12, color=GREY_DARK)
    add_paragraph(tf,
                  "Logit roughly doubles Diabetes recall — most balanced single deployable model.",
                  size=12, color=GREY_DARK, space_before=6)
    add_paragraph(tf,
                  "Only KNN k=5 ever predicts Prediabetes (88 times, 2 correct).",
                  size=12, color=GREY_DARK, space_before=6)

    add_textbox(s, Inches(0.5), Inches(6.7), Inches(7.0), Inches(0.4),
                "Top 5 (rust) align directly with the EDA findings — same five variables.",
                size=12, italic=True, color=GREY_MID)

    add_footer(s, 9, TOTAL)


def slide_10_conclusion():
    s = add_blank()
    add_slide_title(s, "Practical Value and Conclusion", 10, "PRACTICAL VALUE")

    # Left: Rule-out branch (with proportion badge)
    add_rounded_card(s, Inches(0.5), Inches(1.7),
                     Inches(6.1), Inches(3.4),
                     fill=GREY_BG)
    add_textbox(s, Inches(0.85), Inches(1.95), Inches(5.5), Inches(0.4),
                "RULE-OUT BRANCH",
                size=11, bold=True, color=TEAL)
    # Proportion badge top-right of card
    badge_l = add_rounded_card(s, Inches(5.10), Inches(1.92),
                               Inches(1.25), Inches(0.4),
                               fill=TEAL_TINT, border=TEAL, border_w=1.0)
    add_textbox(s, Inches(5.10), Inches(1.97), Inches(1.25), Inches(0.3),
                "~40% of patients", size=9, bold=True, color=TEAL,
                align=PP_ALIGN.CENTER)

    add_textbox(s, Inches(0.85), Inches(2.40), Inches(5.5), Inches(1.0),
                "101,437", size=42, bold=True, color=NAVY)
    add_textbox(s, Inches(0.85), Inches(3.30), Inches(5.5), Inches(0.45),
                "patients with HighBP = 0",
                size=14, color=GREY_DARK)
    add_accent_line(s, Inches(0.85), Inches(3.85), Inches(0.7), color=TEAL)
    add_textbox(s, Inches(0.85), Inches(3.95), Inches(5.5), Inches(1.1),
                "Only 6 % Diabetes rate. Safe to deprioritize from intensive screening.",
                size=14, color=GREY_DARK)

    # Right: High-risk leaf (with proportion badge)
    add_rounded_card(s, Inches(6.8), Inches(1.7),
                     Inches(6.1), Inches(3.4),
                     fill=RUST_TINT, border=RUST, border_w=1.5)
    add_textbox(s, Inches(7.15), Inches(1.95), Inches(5.5), Inches(0.4),
                "HIGH-RISK LEAF",
                size=11, bold=True, color=RUST)
    badge_r = add_rounded_card(s, Inches(11.40), Inches(1.92),
                               Inches(1.25), Inches(0.4),
                               fill=RUST, border=RUST, border_w=1.0)
    add_textbox(s, Inches(11.40), Inches(1.97), Inches(1.25), Inches(0.3),
                "~2% of patients", size=9, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)

    add_textbox(s, Inches(7.15), Inches(2.40), Inches(5.5), Inches(1.0),
                "3,961", size=42, bold=True, color=RUST)
    add_textbox(s, Inches(7.15), Inches(3.30), Inches(5.5), Inches(0.45),
                "BMI ≥ 35  +  HighChol  +  poor GenHlth  +  HighBP",
                size=12, color=GREY_DARK)
    add_accent_line(s, Inches(7.15), Inches(3.85), Inches(0.7), color=RUST)
    add_textbox(s, Inches(7.15), Inches(3.95), Inches(5.5), Inches(1.1),
                "~60 % Diabetes rate — about 5× the population base rate. Tractable population for targeted intervention.",
                size=14, color=GREY_DARK)

    add_rounded_card(s, Inches(0.5), Inches(5.3),
                     Inches(12.4), Inches(1.0), fill=NAVY)
    add_textbox(s, Inches(0.85), Inches(5.4), Inches(11.7), Inches(0.4),
                "HONEST FRAMING", size=11, bold=True, color=AMBER)
    add_textbox(s, Inches(0.85), Inches(5.65), Inches(11.7), Inches(0.6),
                "Accuracy ≈ baseline; class imbalance dominates. Best next step is methodological — class weights, oversampling, or a cost-sensitive threshold — not a different algorithm.",
                size=14, color=WHITE)

    add_textbox(s, Inches(0.5), Inches(6.5), Inches(12.4), Inches(0.4),
                "AI disclosure: generative AI organized + polished structure; data, methods, and analytical decisions are the author's, supported by validated R outputs.",
                size=10, italic=True, color=GREY_MID)

    add_footer(s, 10, TOTAL)


# ----------------------------------------------------------------------- build

slide_1_title()
slide_2_why()
slide_3_dataset()
slide_4_variables()
slide_5_eda()
slide_6_methods()
slide_7_comparison()
slide_8_tree()
slide_9_importance()
slide_10_conclusion()

prs.save(DECK)
print(f"Saved {DECK}")
